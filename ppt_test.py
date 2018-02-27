from pptx.util import Inches
from pptx.util import Pt
from pptx import *
import re


class pptData:
    # title = "P-" + ord + invTi
    # 문헌번호 = apptCtry + docNum + docDt
    WipsOrd = ""  # ord is reserved word in Python
    invTi = ""
    apptCtry = ""
    docNum = ""
    docDt = ""
    applNum = ""
    appt = ""
    invt = ""
    ipc = ""
    cpc = ""
    prList = ""
    myfolderPatMemo = ""
    classificationTag = ""
    ab = ""
    exmpCl = ""
    fmlyList = ""
    ### to be modified later IMAGE PATH
    exmplDrw = ""
    drwList = ""

    def __init__(self, WipsOrd, invTi, apptCtry, docNum, docDt, applNum, appt, invt, ipc, cpc, prList, myfolderPatMemo,
                 classificationTag, ab, exmpCl, fmlyList, exmplDrw, drwList):
        self.WipsOrd = WipsOrd
        self.invTi = invTi
        self.apptCtry = apptCtry
        self.docNum = docNum
        self.docDt = docDt
        self.applNum = applNum  # 출원번호 고치기
        self.appt = appt
        self.invt = invt
        self.ipc = ipc
        self.cpc = cpc
        self.prList = prList
        self.myfolderPatMemo = myfolderPatMemo
        self.classificationTag = classificationTag
        self.ab = ab
        self.exmpCl = exmpCl
        self.fmlyList = fmlyList
        ### to be modified later IMAGE PATH
        self.exmplDrw = exmplDrw
        self.drwList = drwList


# merge cells vertically
def mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx):
    row_count = end_row_idx - start_row_idx + 1
    column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]

    column_cells[0]._tc.set('rowSpan', str(row_count))
    for c in column_cells[1:]:
        c._tc.set('vMerge', '1')


# merge cells horizontally
def mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx):
    col_count = end_col_idx - start_col_idx + 1
    row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
    row_cells[0]._tc.set('gridSpan', str(col_count))
    for c in row_cells[1:]:
        c._tc.set('hMerge', '1')


# the workaround function to merge cells in a table
def mergeCells(table, start_row_idx, end_row_idx, start_col_idx, end_col_idx):
    for col_idx in range(start_col_idx, end_col_idx + 1):
        mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx)
    for row_idx in range(start_row_idx, end_row_idx + 1):
        mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def getData():
    #### Hard Coded part to be modified later
    dataList = []
    WipsOrd = "움직임 탐색시 효율적인 움직임 벡터 추출 방법 및 그 장치"
    invTi = "(Method and Apparatus for effective motion vector decision for motion estimation)"
    apptCtry = "KR"
    docNum = "10-2017-0095793 A "
    docDt = "(2017.08.23)"
    applNum = "10-2017-0102495 (2017.08.11)"
    appt = "한국전자통신연구원"
    invt = "김익균 | 신경선 | 엄낙웅 | 정희범"
    ipc = "H04N-019/573"
    cpc = "H04N-0019/573"
    prList = ""
    myfolderPatMemo = ""
    classificationTag = ""
    ab = "움직임 탐색시 효율적인 움직임 벡터 추출 방법 및 그 장치가 개시된다. 원영상에서 탐색 개시 위치를 결정하여 나선형 움직임 탐색을 수행하는 단계 및 P 픽처 탐색시 서브 샘플링 영상에서의 탐색 수행 여부를 판정하는 단계를 포함하는 움직임 벡터 추출 방법은 서브 샘플링 영상을 이용하면서 나선형움직임탐색, 확장 탬플릿의 복수 병용이라는 방안을 조합한 새로운 계층 나선형 움직임탐색방법인 서브 샘플링탐색에 의한 다수의 움직임벡터후보 검출에 의해 탐색 정도(accuracy)를 개선할 수 있는 효과가 있다."
    exmpCl = "대상 블록에 대한 복수의 움직임 정보 후보를 유도하는 단계; 및상기 복수의 움직임 정보 후보를 기초로 상기 대상 블록의 움직임 정보를 결정하는 단계를 포함하되,상기 복수의 움직임 정보 후보를 유도하는 단계는,기 도출된 제1 움직임 정보 성분 및 제2 움직임 정보 성분에 기초한 차이를 임계치와 비교하는 단계; 및상기 비교 결과에 기초하여 상기 대상 블록의 움직임 정보 후보를 유도하는 단계를 포함하는 움직임 정보 결정 방법."
    fmlyList = "KR10-1677696B1 | KR10-1708905B1 | KR10-1737606B1 | KR10-1769575B1 | KR10-2017-0095793A | US8989269B2"
    ### To be modified later IMAGE PATH
    exmplDrw = "img1.png"
    drwList = "img2.png,img2.png,img2.png,img2.png"
    data1 = pptData(WipsOrd, invTi, apptCtry, docNum, docDt, applNum, appt, invt, ipc, cpc, prList, myfolderPatMemo,
                    classificationTag, ab, exmpCl, fmlyList, exmplDrw, drwList)
    dataList.append(data1)  #########

    WipsOrd = "움직임 탐색시 효율적인 움직임 벡터fective motion or decis decision for m임 탐색시 효율적인 움직임 임 탐색시 효율적인 움직임 fective motion vector decis임 탐색시 효율적인 움직임 ion for m 추출 방법 및 그 장치"
    invTi = "(Method and Apparatus for effective motion vector decision for motion estimation)"
    apptCtry = "KR"
    docNum = "10-2017-0095793 A "
    docDt = "(2017.08.23)"
    applNum = "10-2017-0102495 (2017.08.11)"
    appt = "한국전자통신연구원"
    invt = "김익균 | 신경선 | 엄낙웅 | 정희범"
    ipc = "H04N-019/573"
    cpc = "H04N-0019/573"
    prList = ""
    myfolderPatMemo = ""
    classificationTag = ""
    ab = "움직임 탐색시 효율적인 움직임 벡터 추출 방법 및 그 장치가 개시된다. ㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ원영상에서 탐색 개시 위치를 결정하여 나선형 움직임 탐색을 수행하는 단계 및 P 픽처 탐색시 서브 샘플링 영상에서의 탐색 수행 여부를 판정하는 단계를 포함하는 움직임 벡터 추출 방법은 서브 샘플링 영상을 이용하면서 나선형움직임탐색, 확장 탬플릿의 복수 병용이라는 방안을ㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ 조합한 새로운 계층 나선형 움직임탐색방법인 서브 샘플링탐색에 의한 다수의 움직임벡터후보 검출에 의해 탐색 정도(accuracy)를 개선할 수 있는 효과가 있다."
    exmpCl = "대상 블록에 대한 복수의 움직임 정보 후보를 유도하는 단계; 및상기 복수의 움ㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ임 정보 후보를 기초로 상기 대상 블록의 움직임 정보를 결정하는 단계를 포함하되,상기 복수의 움직임 정보 후보를 유도하는 단계는,기 도출된 제1 움직임 정보 성분 및 제2 움직임 정보 성분에 기초한 차이를 임계치와 비교하는 단계; 및상기 비교 결과에 기초하여 상기 대상 블록의 움직임 정보 후보를 유도하는 단계를 포함하는 움직임 정보 결정 방법."
    fmlyList = "KR10-1677696B1 | KR10-1708905B1 | KR10-1737606B1 | KR10-1769575B1 | KR10-2017-0095793A | US8989269B2"
    ### To be modified later IMAGE PATH
    exmplDrw = "img1.png"
    drwList = "img2.png,img2.png,img2.png,img2.png"
    data2 = pptData(WipsOrd, invTi, apptCtry, docNum, docDt, applNum, appt, invt, ipc, cpc, prList, myfolderPatMemo,
                    classificationTag, ab, exmpCl, fmlyList, exmplDrw, drwList)
    dataList.append(data2)  #########

    WipsOrd = "제목이 길어서 테이블 높이를 바꾸는 경우 제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우제목이 길어서 테이블 높이를 바꾸는 경우"
    invTi = "(In case of long input changes table cell's height In case of long input changes table cell's height In case of long input changes table cell's height In case of long input changes table cell's height In case of long input changes table cell's height In case of long input changes table cell's height )"
    apptCtry = "KR"
    docNum = "10-2017-0095793 A "
    docDt = "(2017.08.23)"
    applNum = "10-2017-0102495 (2017.08.11)"
    appt = "한국전자통신연구원"
    invt = "김익균 | 신경선 | 엄낙웅 | 정희범"
    ipc = "H04N-019/573"
    cpc = "H04N-0019/573"
    prList = ""
    myfolderPatMemo = ""
    classificationTag = ""
    ab = "움직임 탐색시 효율적인 움직임 벡터 추출 방법 및 그 장치가 개시된다. ㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ원영상에서 탐색 개시 위치를 결정하여 나선형 움직임 탐색을 수행하는 단계 및 P 픽처 탐색시 서브 샘플링 영상에서의 탐색 수행 여부를 판정하는 단계를 포함하는 움직임 벡터 추출 방법은 서브 샘플링 영상을 이용하면서 나선형움직임탐색, 확장 탬플릿의 복수 병용이라는 방안을ㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ 조합한 새로운 계층 나선형 움직임탐색방법인 서브 샘플링탐색에 의한 다수의 움직임벡터후보 검출에 의해 탐색 정도(accuracy)를 개선할 수 있는 효과가 있다."
    exmpCl = "대상 블록에 대한 복수의 움직임 정보 후보를 유도하는 단계; 및상기 복수의 움ㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹㄻㄴㅇㄻㅁㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄻㄴㅇㄹ임 정보 후보를 기초로 상기 대상 블록의 움직임 정보를 결정하는 단계를 포함하되,상기 복수의 움직임 정보 후보를 유도하는 단계는,기 도출된 제1 움직임 정보 성분 및 제2 움직임 정보 성분에 기초한 차이를 임계치와 비교하는 단계; 및상기 비교 결과에 기초하여 상기 대상 블록의 움직임 정보 후보를 유도하는 단계를 포함하는 움직임 정보 결정 방법."
    fmlyList = "KR10-1677696B1 | KR10-1708905B1 | KR10-1737606B1 | KR10-1769575B1 | KR10-2017-0095793A | US8989269B2"
    ### To be modified later IMAGE PATH
    exmplDrw = "img1.png"
    drwList = "img2.png,img2.png,img2.png,img2.png"
    data3 = pptData(WipsOrd, invTi, apptCtry, docNum, docDt, applNum, appt, invt, ipc, cpc, prList, myfolderPatMemo,
                    classificationTag, ab, exmpCl, fmlyList, exmplDrw, drwList)
    dataList.append(data3)  #########

    return dataList


def lineCounter(data):
    lineCount = 0
    hangul = re.compile('[ㄱ-ㅣ가-힣]')
    eng = re.compile('[A-Za-z0-9]')
    short = re.compile(',|\.|;|:|\'|!|`| ')

    titleStr = "P-." + data.WipsOrd + data.invTi
    totLen = len(titleStr)
    hanLen = len(hangul.findall(titleStr))
    engLen = len(eng.findall(titleStr))
    shortLen = len(short.findall(titleStr))
    elseLen = totLen - (hanLen + engLen + shortLen)
    # print("tot{} han{} eng{} short{} else{}".format(totLen, hanLen, engLen, shortLen, elseLen))

    hanSize = 0.0135135
    engSize = 0.00704225
    shortSize = 0.003571427
    elseSize = 0.00409836

    hanLine = hanLen * hanSize
    engLine = engLen * engSize
    shortLine = shortLen * shortSize
    elseLine = elseLen * elseSize

    # print("han:{} eng:{} short:{} else:{}".format(hanLine, engLine, shortLine, elseLine))

    num = hanLine + engLine + shortLine + elseLine
    # print(num)
    while num > 1:
        num = num - 1
        lineCount = lineCount + 1

    return lineCount


def createPpt(dataList):  # copy_slide
    prs = Presentation()
    i = 1
    emptySlide = prs.slide_layouts[6]
    rows = 11
    cols = 4
    left = top = Inches(0.2)
    height = Inches(2)
    width = Inches(9.5)

    for data in dataList:
        slide = prs.slides.add_slide(emptySlide)
        shapes = slide.shapes
        table = shapes.add_table(rows, cols, left, top, width, height).table
        fillTmpl(table, data)

        lines = lineCounter(data)
        print(lines)
        ileft, itop = 1.5, (4.3 + lines * 0.125)
        slide.shapes.add_picture("img1.png", Inches(ileft), Inches(itop))

        ####################lineCount(data)
        itop = itop + 1.45
        for i in range(0, 4):
            slide.shapes.add_picture("img2.png", Inches(ileft), Inches(itop))
            ileft = ileft + 1.5

    prs.save('test2.pptx')

def fillTmpl(table, data):
    table.rows[6].height = Inches(1.25)
    table.rows[7].height = Inches(1.25)
    table.rows[8].height = Inches(1.45)
    table.rows[9].height = Inches(1.45)
    table.columns[0].width = Inches(1.25)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(1.25)
    table.columns[3].width = Inches(3.5)

    mergeCellsHorizontally(table=table, row_idx=0, start_col_idx=0, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=5, start_col_idx=1, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=6, start_col_idx=1, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=7, start_col_idx=1, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=8, start_col_idx=1, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=9, start_col_idx=1, end_col_idx=3)
    mergeCellsHorizontally(table=table, row_idx=10, start_col_idx=1, end_col_idx=3)

    table.cell(1, 0).text = "문헌번호"
    table.cell(1, 2).text = "출원번호"
    table.cell(2, 0).text = "출원인"
    table.cell(2, 2).text = "발명자"
    table.cell(3, 0).text = "Current IPC(Main)"
    table.cell(3, 2).text = "Current CPC(Main)"
    table.cell(4, 0).text = "우선권번호"
    table.cell(4, 2).text = "문헌메모"
    table.cell(5, 0).text = "사용자태그"
    table.cell(6, 0).text = "요약"
    table.cell(7, 0).text = "대표청구항"
    table.cell(8, 0).text = "대표도면"
    table.cell(9, 0).text = "개별도면"
    table.cell(10, 0).text = "Wips패밀리"

    # table.cell(0,0).text = "P-." + data.WipsOrd + data.invTi
    table.cell(0, 0).text = data.WipsOrd + data.invTi
    table.cell(1, 1).text = data.apptCtry + data.docNum + data.docDt
    table.cell(1, 3).text = data.applNum
    table.cell(2, 1).text = data.appt
    table.cell(2, 3).text = data.invt
    table.cell(3, 1).text = data.ipc
    table.cell(3, 3).text = data.cpc
    table.cell(4, 1).text = data.prList
    table.cell(4, 3).text = data.myfolderPatMemo
    table.cell(5, 1).text = data.classificationTag
    table.cell(6, 1).text = data.ab
    table.cell(7, 1).text = data.exmpCl
    # 대표도면
    # 개별도면
    table.cell(10, 1).text = data.fmlyList

    for cell in iter_cells(table):
        cell.text_frame.paragraphs[0].font.size = Pt(9)

    # 테스트 111111
    # 테스트2222222222
    return table

if __name__ == "__main__":
    dataList = getData()
    createPpt(dataList)
