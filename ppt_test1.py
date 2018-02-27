from pptx.util import Inches
from pptx.util import Pt
from pptx import *
import re

def createPpt():  # copy_slide
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    #left = top = width = height = Inches(2)
    left = Inches(1)
    top = Inches(1)
    width = Inches(4)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)

    tf = txBox.text_frame
    tf.text = "This is text inside a textbox"

    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)

    left = Inches(1)
    top = Inches(2)
    width = Inches(4)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Just an example"
    font = run.font
    font.size = Pt(30)

    prs2 = Presentation()
    title_slide_layout = prs2.slide_layouts[0]
    slide = prs2.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')

    print("out")
    print (Presentation())

    return ""

if __name__ == "__main__":
    createPpt()
